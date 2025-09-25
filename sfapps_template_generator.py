"""
sfapps_template_generator.py
=================================

This module provides a high-level helper for constructing
PowerPoint presentations from a pre-built template.  Unlike
``sfapps_presentation_generator.py`` – which rebuilds each slide from
scratch – this module edits an existing PPTX file in place.  It
maintains all decorative artwork, backgrounds and layout elements
defined in the original template and only replaces the fields the
caller cares about: the cover page title, the per-application name,
publisher and logo, and the closing call-to-action.

The workflow is as follows:

* The caller specifies a topic (industry/category), a list of
  AppExchange URLs and a URL for the final slide.  Optionally a
  dictionary of overrides can provide custom names, publishers and
  logo images when automatic extraction fails or should be
  customised.
* The template file is opened using ``python-pptx``.  The template
  shipped with this assignment contains a cover slide, ten
  programme slides (#1 through #10) and a closing slide.  If more
  than ten applications are supplied the second programme slide is
  duplicated until enough placeholders exist.
* Each programme slide is updated: the number is set to ``#1``,
  ``#2`` etc., the app name and publisher text is replaced, and the
  logo image is swapped.  Decorative artwork is left untouched.
* Any unused programme slides are removed so the final deck only
  contains the cover, the required programme slides and the
  closing slide.
* The cover slide is updated by replacing the ``$industry`` token
  with the supplied topic.  The closing slide is similarly updated
  and the logo on that slide is given a hyperlink to the final URL.
* The result is written to ``output_pptx``.  If LibreOffice is
  available a PDF is also produced at ``output_pdf``.

The module exposes a single public entry point:

    create_presentation_from_template(topic, links, final_url,
                                      template_path, output_pptx,
                                      output_pdf=None,
                                      app_overrides=None)

Requirements
------------

* python-pptx
* requests
* beautifulsoup4
* Pillow (for image scaling)
* LibreOffice in headless mode (optional, for PDF conversion)
"""

import os
import re
import subprocess
from copy import deepcopy
from dataclasses import dataclass, field
from io import BytesIO
from typing import Dict, List, Optional, Tuple

import requests
from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt


@dataclass
class AppMetadata:
    """Container for app details extracted from an AppExchange listing."""

    url: str
    name: str
    developer: str
    logo_bytes: bytes
    logo_mime: str


def _extract_from_html(html: str) -> Tuple[Optional[str], Optional[str], Optional[str]]:
    """
    Given the HTML body of an AppExchange listing this helper will try
    to extract the application name, the developer/publisher and a
    logo URL.  CSS selectors are prioritized, falling back to JSON and 
    OpenGraph metadata.

    Parameters
    ----------
    html: str
        Raw HTML string from the AppExchange listing page.

    Returns
    -------
    (name, developer, logo_url): Tuple of three strings or ``None`` if
    a field cannot be determined.
    """
    soup = BeautifulSoup(html, 'html.parser')
    name = None
    dev = None
    logo = None
    
    # Try CSS selectors first (most reliable for AppExchange)
    # Correct selectors for title (based on HTML structure)
    name_selectors = [
        'h1[type="style"]',  # Visible in HTML on the right
        '.listing-title h1',
        'h1',
        '[data-testid="listing-title"]'
    ]
    for selector in name_selectors:
        element = soup.select_one(selector)
        if element:
            text = element.get_text().strip()
            if text:  # Make sure text is not empty
                name = text
                print(f"Found title via selector '{selector}': {name}")
                break

    # Correct selectors for developer (based on HTML structure)
    dev_selectors = [
        'p[type="style"]',  # Visible in HTML on the right - "By TaskRay"
        '.listing-title p',
        'p',
        '[data-testid="listing-publisher"]'
    ]
    for selector in dev_selectors:
        element = soup.select_one(selector)
        if element:
            dev_text = element.get_text().strip()
            if dev_text:  # Make sure text is not empty
                # Remove "By " prefix if present
                if dev_text.lower().startswith('by '):
                    dev = dev_text[3:].strip()
                else:
                    dev = dev_text
                print(f"Found developer via selector '{selector}': {dev}")
                break

    # Correct selectors for logo (based on HTML structure)
    logo_selectors = [
        'img.ads-image',  # Exactly as shown in HTML on the right
        '.ads-image',
        '.listing-logo img',
        '.summary img',
        'img[class*="ads-image"]'
    ]
    for selector in logo_selectors:
        element = soup.select_one(selector)
        if element:
            # Try different attributes for image URL
            logo = element.get('src') or element.get('data-src') or element.get('data-original') or element.get('data-lazy')
            if logo:
                print(f"Found logo via selector '{selector}': {logo}")
                break
    
    # Try to extract from JSON script tags only if CSS failed
    if not all([name, dev, logo]):
        script_tags = soup.find_all('script', type='application/json')
        for script in script_tags:
            try:
                import json
                data = json.loads(script.get_text())
                # Try to find app data in JSON structure
                if isinstance(data, dict):
                    # Look for common patterns in JSON data
                    for key, value in data.items():
                        if isinstance(value, dict):
                            if 'name' in value and 'developer' in value:
                                name = name or value.get('name')
                                dev = dev or value.get('developer')
                            elif 'title' in value and 'publisher' in value:
                                name = name or value.get('title')
                                dev = dev or value.get('publisher')
            except:
                continue
    
    # Final fallback to OpenGraph metadata
    if not name:
        og_title = soup.find('meta', property='og:title')
        if og_title and og_title.get('content'):
            title_content = og_title['content']
            # Remove common suffixes like "| Salesforce AppExchange"
            if '|' in title_content:
                name = title_content.split('|')[0].strip()
            else:
                name = title_content.strip()
    
    if not logo:
        og_image = soup.find('meta', property='og:image')
        if og_image and og_image.get('content'):
            logo = og_image['content']
    
    if not dev:
        # Look for Twitter metadata
        twitter_data1 = soup.find('meta', attrs={'name': 'twitter:data1'})
        if twitter_data1 and twitter_data1.get('content'):
            dev = twitter_data1['content'].strip()
        else:
            # Look for any span/text containing "By"
            by_elements = soup.find_all(string=re.compile(r'By\s+', re.IGNORECASE))
            for by_text in by_elements:
                if by_text.strip():
                    dev = by_text.replace('By', '').strip()
                    break
    
    return name, dev, logo


def fetch_app_metadata(url: str, timeout: int = 20) -> Optional[AppMetadata]:
    """
    Retrieve metadata for an AppExchange listing using modern Selenium parser.
    If extraction or download fails, ``None`` is returned.

    Parameters
    ----------
    url: str
        URL of the AppExchange listing.
    timeout: int, optional
        Maximum number of seconds to wait for HTTP requests.

    Returns
    -------
    AppMetadata or None
        ``AppMetadata`` containing the name, developer and logo bytes
        if successful, otherwise ``None``.
    """
    # Import modern Selenium parser
    try:
        from appexchange_parser import parse_appexchange_improved
        print(f"🔄 Using modern Selenium parser for {url}")
        
        # Use modern parser with Shadow DOM support
        result = parse_appexchange_improved(url)
        
        if not result or not result.get('success'):
            print(f"❌ Parser could not extract data from {url}")
            return None
            
        name = result.get('name', 'Unknown App')
        developer = result.get('developer', 'Unknown Developer')
        logo_url = result.get('logo_url')
        
        print(f"✅ Selenium parser extracted data:")
        print(f"   Name: {name}")
        print(f"   Developer: {developer}")
        print(f"   Logo URL: {logo_url}")
        
        # Download logo
        logo_bytes = b''
        logo_mime = 'image/png'
        
        if logo_url:
            try:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 '
                                  '(KHTML, like Gecko) Chrome/120.0 Safari/537.36'
                }
                img_resp = requests.get(logo_url, headers=headers, timeout=timeout)
                img_resp.raise_for_status()
                logo_bytes = img_resp.content
                logo_mime = img_resp.headers.get('Content-Type', 'image/png')
                print(f"✅ Logo downloaded: {len(logo_bytes)} bytes, MIME: {logo_mime}")
            except Exception as e:
                print(f"⚠️ Logo download error: {e}")
                logo_bytes = b''
        
        return AppMetadata(url=url, name=name, developer=developer, logo_bytes=logo_bytes, logo_mime=logo_mime)
        
    except ImportError:
        print(f"❌ Selenium parser unavailable, using fallback for {url}")
        # Fallback to old HTML parser only if Selenium unavailable
        headers = {
            'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 '
                          '(KHTML, like Gecko) Chrome/120.0 Safari/537.36'
        }
        try:
            resp = requests.get(url, headers=headers, timeout=timeout)
            resp.raise_for_status()
        except Exception:
            return None
        name, dev, logo_url = _extract_from_html(resp.text)
        if not name or not dev or not logo_url:
            return None
        # Fetch logo
        try:
            img_resp = requests.get(logo_url, headers=headers, timeout=timeout)
            img_resp.raise_for_status()
            logo_bytes = img_resp.content
            logo_mime = img_resp.headers.get('Content-Type', 'image/png')
        except Exception:
            return None
        return AppMetadata(url=url, name=name, developer=dev, logo_bytes=logo_bytes, logo_mime=logo_mime)
    except Exception as e:
        print(f"❌ Error in fetch_app_metadata: {e}")
        return None


def _remove_comments_from_slides(prs: Presentation, slide_indices: List[int]) -> None:
    """
    Remove all comments from specified slides in the presentation.
    
    Parameters
    ----------
    prs: Presentation
        The presentation object to operate on.
    slide_indices: List[int]
        Zero-based indices of slides to remove comments from.
    """
    print(f"🗑️ Removing comments from slides: {[i+1 for i in slide_indices]}")
    
    for slide_idx in slide_indices:
        if slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            try:
                slide_part = slide.part
                
                # Find all relationships to comments (multiple search variants)
                comment_rels = []
                
                print(f"   🔍 Analyzing slide {slide_idx + 1}, found relationships: {len(slide_part.rels)}")
                
                for rel_id, rel in slide_part.rels.items():
                    rel_type = getattr(rel, 'reltype', 'unknown')
                    print(f"     - {rel_id}: {rel_type}")
                    
                    # Look for different comment variants
                    if (hasattr(rel, 'reltype') and 
                        ('comment' in rel_type.lower() or 
                         'comments' in rel_type.lower() or
                         rel_type.endswith('/comments'))):
                        comment_rels.append(rel_id)
                        print(f"     ✅ Found comment: {rel_id} ({rel_type})")
                
                # Remove comment relationships
                for rel_id in comment_rels:
                    try:
                        comment_part = slide_part.rels[rel_id].target_part
                        
                        slide_part.drop_rel(rel_id)
                        
                        try:
                            if hasattr(prs.part, 'package'):
                                package = prs.part.package
                                if hasattr(package, '_parts') and comment_part.partname in package._parts:
                                    del package._parts[comment_part.partname]
                        except Exception as pkg_e:
                            print(f"   ⚠️ Could not remove from package: {pkg_e}")
                        
                        print(f"   ✅ Removed comment {rel_id} from slide {slide_idx + 1}")
                        
                    except Exception as e:
                        print(f"   ⚠️ Error removing comment {rel_id}: {e}")
                
                if not comment_rels:
                    print(f"   ℹ️ No comments found on slide {slide_idx + 1}")
                else:
                    print(f"   ✅ Processed {len(comment_rels)} comments on slide {slide_idx + 1}")
                    
            except Exception as e:
                print(f"   ⚠️ Error removing comments from slide {slide_idx + 1}: {e}")
        else:
            print(f"   ⚠️ Slide {slide_idx + 1} not found (total slides: {len(prs.slides)})")


def _clone_slide(prs: Presentation, index: int) -> None:
    """
    Clone the slide at position ``index`` and append the clone to the
    end of the slide collection.  All shapes on the source slide are
    deep–copied; group shapes are not present in the provided
    template so no special casing is necessary.  The new slide uses
    the same layout as the source slide.

    Parameters
    ----------
    prs: Presentation
        The presentation object to operate on.
    index: int
        Zero based index of the slide to clone.
    """
    source = prs.slides[index]
    layout = prs.slide_layouts[0]  # Only one layout defined in template
    new_slide = prs.slides.add_slide(layout)
    for shape in source.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape.element))


def _remove_slide(prs: Presentation, index: int) -> None:
    """
    Remove the slide at the specified index from the presentation.

    This helper updates the slide id list and drops the corresponding
    relationship from the presentation part.  Without both steps the
    PPTX would become corrupt.

    Parameters
    ----------
    prs: Presentation
        The presentation from which to remove a slide.
    index: int
        Zero based slide index to remove.
    """
    from pptx.oxml.ns import qn
    slide = prs.slides[index]
    slide_id = slide.slide_id
    sldIdLst = prs.slides._sldIdLst
    relId = None
    for sldId in list(sldIdLst):
        if int(sldId.get('id')) == slide_id:
            relId = sldId.get(qn('r:id'))
            sldIdLst.remove(sldId)
            break
    if relId:
        prs.part.drop_rel(relId)


def _calculate_text_width(text, font_size, font_name='Poppins', bold=False):
    """Calculate approximate text width in points"""
    # Character width ratios for different fonts (approximate values)
    font_ratios = {
        'Poppins': 0.6 if not bold else 0.65,
        'Arial': 0.55 if not bold else 0.6,
        'Times New Roman': 0.5 if not bold else 0.55,
    }
    
    # Use Poppins ratio as default if font not found
    ratio = font_ratios.get(font_name, 0.6)
    
    # Calculate base width: chars * font_size * ratio
    base_width = len(text) * font_size * ratio
    
    # Add padding (about 20% for comfortable spacing)
    padding = base_width * 0.2
    
    return base_width + padding


def _remove_developer_background(slide, text_left, text_top, text_height):
    """
    Find and remove the blue background shape behind developer text
    
    Parameters
    ----------
    slide: pptx.slide.Slide
        The slide containing the shapes
    text_left: int
        Left position of the text field (in EMU)
    text_top: int
        Top position of the text field (in EMU)
    text_height: int
        Height of the text field (in EMU)
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    
    print(f"   🗑️ Searching for blue background to remove...")
    print(f"      Text position: left={text_left/914400:.1f}in, top={text_top/914400:.1f}in")
    
    tolerance = Pt(100).emu  # Increased tolerance for position matching
    shapes_to_remove = []
    
    for idx, shape in enumerate(slide.shapes):
        # Skip text frames
        if shape.has_text_frame:
            continue
            
        print(f"      Checking shape [{idx}]:")
        print(f"        Position: left={shape.left/914400:.1f}in, top={shape.top/914400:.1f}in")
        print(f"        Size: {shape.width/914400:.1f}in x {shape.height/914400:.1f}in")
        
        # Check if shape is close to text position (likely background)
        left_diff = abs(shape.left - text_left)
        top_diff = abs(shape.top - text_top)
        height_diff = abs(shape.height - text_height)
        
        print(f"        Position differences: left={left_diff/914400:.1f}in, top={top_diff/914400:.1f}in, height={height_diff/914400:.1f}in")
        
        # More relaxed matching - check if it's near the text area
        if left_diff < tolerance and top_diff < tolerance:
            print(f"        ✅ Shape is near developer text area")
            
            # Check if it has any fill color (not just blue)
            has_fill = False
            try:
                if hasattr(shape, 'fill'):
                    if hasattr(shape.fill, 'type') and shape.fill.type is not None:
                        has_fill = True
                        print(f"        Shape has fill type: {shape.fill.type}")
                        
                        if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                            r, g, b = shape.fill.fore_color.rgb.r, shape.fill.fore_color.rgb.g, shape.fill.fore_color.rgb.b
                            print(f"        Fill color: RGB({r}, {g}, {b})")
                            
                            # Check for any blue-ish colors (light blue, cyan, etc)
                            if (b > 150 and b > r and b > g) or \
                               (g > 150 and b > 150) or \
                               (150 <= r <= 255 and 180 <= g <= 255 and 200 <= b <= 255):
                                print(f"        ✅ Identified as colored background - marking for removal")
                                shapes_to_remove.append(shape)
                                continue
                
                # If shape is near text but we can't determine color, remove it anyway
                if has_fill or not shape.has_text_frame:
                    print(f"        🔄 Shape near developer text without text - assuming background")
                    shapes_to_remove.append(shape)
                    
            except Exception as e:
                print(f"        ⚠️ Could not check fill properties: {e}")
                # If we can't check properties but position matches, assume it's the background
                print(f"        🔄 Assuming background based on position - marking for removal")
                shapes_to_remove.append(shape)
    
    # Remove the identified shapes
    if shapes_to_remove:
        print(f"      Found {len(shapes_to_remove)} shapes to remove")
        for i, shape in enumerate(shapes_to_remove):
            try:
                slide.shapes._spTree.remove(shape._element)
                print(f"        ✅ Background shape {i+1} removed successfully")
            except Exception as e:
                print(f"        ❌ Failed to remove background shape {i+1}: {e}")
    else:
        print(f"      ❌ No background shapes found near developer text")


def _update_developer_background(slide, text_left, text_top, text_height, target_width):
    """
    Find and update the blue background shape behind developer text
    
    Parameters
    ----------
    slide: pptx.slide.Slide
        The slide containing the shapes
    text_left: int
        Left position of the text field (in EMU)
    text_top: int
        Top position of the text field (in EMU)
    text_height: int
        Height of the text field (in EMU)
    target_width: float
        Target width in points
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    
    print(f"   🔍 Searching for blue background near developer text...")
    print(f"      Text position: left={text_left/914400:.1f}in, top={text_top/914400:.1f}in")
    
    target_width_emu = Pt(target_width).emu
    tolerance = Pt(50).emu  # 50pt tolerance for position matching
    
    for idx, shape in enumerate(slide.shapes):
        # Skip text frames
        if shape.has_text_frame:
            continue
            
        # Check if shape is close to text position (likely background)
        left_diff = abs(shape.left - text_left)
        top_diff = abs(shape.top - text_top)
        height_diff = abs(shape.height - text_height)
        
        # Must be close in position and similar height
        if left_diff < tolerance and top_diff < tolerance and height_diff < tolerance:
            print(f"      Found potential background shape [{idx}]:")
            print(f"        Position: left={shape.left/914400:.1f}in, top={shape.top/914400:.1f}in")
            print(f"        Size: {shape.width/914400:.1f}in x {shape.height/914400:.1f}in")
            
            # Check if it has a blue-ish fill color
            try:
                if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                    fill_color = shape.fill.fore_color
                    if hasattr(fill_color, 'rgb'):
                        r, g, b = fill_color.rgb.r, fill_color.rgb.g, fill_color.rgb.b
                        print(f"        Fill color: RGB({r}, {g}, {b})")
                        
                        # Check for light blue colors (typical background colors)
                        if (150 <= r <= 255 and 180 <= g <= 255 and 200 <= b <= 255) or \
                           (0 <= r <= 100 and 180 <= g <= 255 and 220 <= b <= 255):
                            print(f"        ✅ Identified as blue background - updating width")
                            shape.width = target_width_emu
                            print(f"        📏 Background width updated to: {target_width:.1f}pt")
                            return
            except Exception as e:
                print(f"        ⚠️ Could not check fill color: {e}")
                # If we can't check color but position matches, assume it's the background
                print(f"        🔄 Assuming background based on position - updating width")
                shape.width = target_width_emu
                print(f"        📏 Background width updated to: {target_width:.1f}pt")
                return
    
    print(f"      ❌ No blue background shape found near developer text")


def _find_logo_shape(slide) -> Optional[int]:
    """
    Given a slide, attempt to identify the picture shape that contains
    the application logo.  The heuristic is to pick the picture
    shape whose width and height are both between 1 and 4 inches and
    which has the largest area among such candidates.  Decorative
    icons on the template are much smaller or much larger, so this
    rule works reliably on the provided design.

    Parameters
    ----------
    slide: pptx.slide.Slide
        A slide object to inspect.

    Returns
    -------
    int or None
        The index of the candidate shape or ``None`` if none match.
    """
    print(f"🔍 Searching for logo among {len(slide.shapes)} shapes on slide:")
    candidates: List[Tuple[float, int]] = []
    
    for idx, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w = shape.width / 914400.0  # Convert to inches
            h = shape.height / 914400.0  # Convert to inches
            area = w * h
            print(f"   Shape [{idx}]: {w:.2f}\" x {h:.2f}\" (area: {area:.3f})")
            
            if 1.0 < w < 4.0 and 1.0 < h < 4.0:
                candidates.append((area, idx))
                print(f"     ✅ Candidate for logo")
            else:
                print(f"     ❌ Not suitable (size out of range 1-4 inches)")
        else:
            shape_type_name = str(shape.shape_type).split('.')[-1] if hasattr(shape.shape_type, 'name') else str(shape.shape_type)
            print(f"   Shape [{idx}]: {shape_type_name} (not an image)")
    
    if not candidates:
        print("❌ No suitable shapes found for logo")
        return None
        
    candidates.sort(reverse=True)  # Largest area first
    selected_idx = candidates[0][1]
    selected_area = candidates[0][0]
    print(f"✅ Selected shape [{selected_idx}] with area {selected_area:.3f}")
    return selected_idx


def _update_slide_fields(slide, app: AppMetadata, number: int) -> None:
    """
    Replace the number, name, developer and logo on a single
    programme slide.  The function searches for the first text
    shape containing a hash ('#') and replaces its text with ``#n``.
    It then searches for the first text shape beginning with "By"
    (case–insensitive) and replaces it with the developer string.  Any
    remaining text shape containing the original app name from the
    template is replaced with ``app.name``.  Finally the logo
    placeholder picture is updated by replacing its underlying image
    part.

    Parameters
    ----------
    slide: pptx.slide.Slide
        The slide to operate on.
    app: AppMetadata
        Data for the application to fill in.
    number: int
        One–based sequence number to display on the slide.
    """
    print(f"\n🎯 Updating slide #{number}")
    print(f"   Application: {app.name}")
    print(f"   Developer: {app.developer}")
    print(f"   Logo bytes: {len(app.logo_bytes) if app.logo_bytes else 0} bytes")
    print(f"   Logo MIME: {getattr(app, 'logo_mime', 'not specified')}")
    
    # Update text shapes
    replaced_name = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        if '#' in text:
            # Normalize to a single number with leading space as in the template
            shape.text = f" #{number}"
            # Formatting for number: font Poppins, bold, 40pt, color #ffffff
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Poppins'
                    run.font.bold = True
                    run.font.size = Pt(40)
                    run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
            # Vertical alignment to middle
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            continue
        lowered = text.strip().lower()
        if lowered.startswith('by '):
            shape.text = f"{app.developer}"
            # Formatting for developer: font Poppins, 27pt, left align, color #3cc0ff
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = 'Poppins'
                    run.font.size = Pt(27)
                    run.font.color.rgb = RGBColor(0x3c, 0xc0, 0xff)
            # Vertical alignment to middle
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Calculate optimal width based on text content
            calculated_width = _calculate_text_width(app.developer, 27, 'Poppins', bold=False)
            min_width_pt = 150  # Minimum width for visual consistency
            max_width_pt = 400  # Maximum width to prevent overly wide fields
            
            # Use calculated width but respect min/max boundaries
            optimal_width = max(min_width_pt, min(calculated_width, max_width_pt))
            
            # Store developer text position for background adaptation
            dev_text_left = shape.left
            dev_text_top = shape.top
            dev_text_height = shape.height
            
            # Update text field width
            shape.width = Pt(optimal_width)
            
            # Remove blue background shape behind developer text
            _remove_developer_background(slide, dev_text_left, dev_text_top, dev_text_height)
            
            print(f"   📏 Developer field sizing:")
            print(f"      Text: '{app.developer}' ({len(app.developer)} chars)")
            print(f"      Calculated width: {calculated_width:.1f}pt")
            print(f"      Applied width: {optimal_width:.1f}pt")
            continue
        # Replace the template app name – only the first occurrence
        if not replaced_name and text.strip():
            # If the text originally came from the template it will
            # match one of the placeholder names; simply replace it.
            shape.text = app.name
            # Formatting for name: font Poppins, bold, 40pt, left align, #163560
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = 'Poppins'
                    run.font.bold = True
                    run.font.size = Pt(40)
                    run.font.color.rgb = RGBColor(0x16, 0x35, 0x60)
            # Vertical alignment to middle
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            replaced_name = True
            continue
    # Update logo image
    idx = _find_logo_shape(slide)
    print(f"🔍 Updating logo for {app.name}")
    print(f"   Logo shape index: {idx}")
    print(f"   Logo_bytes size: {len(app.logo_bytes) if app.logo_bytes else 0} bytes")
    print(f"   MIME type: {getattr(app, 'logo_mime', 'not specified')}")
    
    if idx is not None:
        pic_shape = slide.shapes[idx]
        target_width = Pt(207)
        target_height = Pt(161)
        
        if not app.logo_bytes:
            return
            
        # Acquire the relationship id pointing to the image
        rId = pic_shape._element.blip_rId
        image_part = slide.part.related_part(rId)
        print(f"   Relationship ID: {rId}")
        
        # Load image into PIL to scale it down if necessary
        try:
            print("   Loading image into PIL...")
            with Image.open(BytesIO(app.logo_bytes)) as img:
                print(f"   Original image size: {img.size}")
                print(f"   Image format: {img.format}")

                # Size in pixels for 207x161 pt at 96 DPI
                target_w_px = int(207 * 96 / 72)  # ~276 px
                target_h_px = int(161 * 96 / 72)  # ~215 px
                print(f"   Target size for logo: {target_w_px} x {target_h_px} px")

                # Resize while preserving aspect ratio within target bounds
                w, h = img.size
                ratio = min(target_w_px / w, target_h_px / h)
                print(f"   Scaling factor: {ratio:.3f}")

                new_size = (int(w * ratio), int(h * ratio))
                img = img.resize(new_size, Image.LANCZOS)
                print(f"   Image resized to: {new_size}")

                buf = BytesIO()
                img.save(buf, format='PNG')
                new_bytes = buf.getvalue()
                print(f"   Final PNG size: {len(new_bytes)} bytes")

                # Set shape size in PowerPoint, preserving aspect ratio
                # Calculate final sizes in pt for PowerPoint
                final_width_pt = new_size[0] * 72 / 96
                final_height_pt = new_size[1] * 72 / 96
                
                pic_shape.width = Pt(final_width_pt)
                pic_shape.height = Pt(final_height_pt)
                print(f"   Updated shape size: {final_width_pt:.1f}pt x {final_height_pt:.1f}pt")
                
        except Exception as e:
            print(f"❌ Image processing error: {e}")
            print(f"   Using original bytes ({len(app.logo_bytes)} bytes)")
            # If resizing fails, fall back to original bytes, but still set target size
            pic_shape.width = target_width
            pic_shape.height = target_height
            new_bytes = app.logo_bytes
            
        # Overwrite the underlying image part
        print("   Updating image part in presentation...")
        image_part._blob = new_bytes
        print("✅ Logo successfully updated")
    else:
        print("❌ Logo shape not found on slide")


def _update_cover_slide(slide, topic: str) -> None:
    """
    Replace the ``$industry`` token on the cover slide with the
    provided topic string.  Any text shape containing the token will
    have it substituted without altering other characters.

    Parameters
    ----------
    slide: pptx.slide.Slide
        The cover slide to modify.
    topic: str
        The replacement topic string.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if '$industry' in shape.text:
            # Clear existing text
            shape.text = ""
            
            # Create multi-line text with different colors
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            
            # "Best Apps for " - color #163560 (dark blue)
            run1 = paragraph.add_run()
            run1.text = "Best Apps for "
            run1.font.name = 'Poppins'
            run1.font.bold = True
            run1.font.size = Pt(59)
            run1.font.color.rgb = RGBColor(0x16, 0x35, 0x60)
            
            # "{topic}" - color #3cc0ff (light blue)
            run2 = paragraph.add_run()
            run2.text = topic
            run2.font.name = 'Poppins'
            run2.font.bold = True
            run2.font.size = Pt(59)
            run2.font.color.rgb = RGBColor(0x3c, 0xc0, 0xff)
            
            # Add new line for second part
            run3 = paragraph.add_run()
            run3.text = "\n"
            
            # "Available on " - color #163560 (dark blue)
            run4 = paragraph.add_run()
            run4.text = "Available on "
            run4.font.name = 'Poppins'
            run4.font.bold = True
            run4.font.size = Pt(59)
            run4.font.color.rgb = RGBColor(0x16, 0x35, 0x60)
            
            # Vertical alignment to middle
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            print(f"✅ Updated cover slide with topic: 'Best Apps for {topic}\\nAvailable on AppExchange'")


def _update_closing_slide(slide, topic: str, final_url: str) -> None:
    """
    Update the closing slide with the topic and assign a hyperlink to
    the SFApps button/logo.  The template includes two text shapes
    which both contain ``$industry``; these are updated.  A picture
    shape bearing the SFApps logo is given a click hyperlink to
    ``final_url``.  If no such picture can be found the hyperlink
    assignment is silently skipped.

    Parameters
    ----------
    slide: pptx.slide.Slide
        The closing slide.
    topic: str
        The industry/category string.
    final_url: str
        URL to assign to the clickable logo.
    """
    # Replace $industry in text with complex formatting
    for shape in slide.shapes:
        if shape.has_text_frame:
            if '$industry' in shape.text:
                # Create text "Apps for {topic} at"
                full_text = f"Apps for {topic} at"

                # Clear existing text
                shape.text = ""

                # Apply formatting for different parts of the text
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.RIGHT

                # "Apps for " - color #163560
                run1 = paragraph.add_run()
                run1.text = "Apps for "
                run1.font.name = 'Poppins'
                run1.font.bold = True
                run1.font.size = Pt(59)
                run1.font.color.rgb = RGBColor(0x16, 0x35, 0x60)

                # "{topic}" - color #3cc0ff
                run2 = paragraph.add_run()
                run2.text = topic
                run2.font.name = 'Poppins'
                run2.font.bold = True
                run2.font.size = Pt(59)
                run2.font.color.rgb = RGBColor(0x3c, 0xc0, 0xff)

                # " at" - color #163560
                run3 = paragraph.add_run()
                run3.text = " at"
                run3.font.name = 'Poppins'
                run3.font.bold = True
                run3.font.size = Pt(59)
                run3.font.color.rgb = RGBColor(0x16, 0x35, 0x60)

                # Vertical alignment to middle
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                print(f"✅ Updated closing slide text: 'Apps for {topic} at'")
    # Assign hyperlink to picture containing SFApps logo; heuristic is
    # to pick the image with a long width and small height (the pill
    # shaped button) – this is picture index 3 in the template.
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Calculate aspect ratio; the SFApps pill button is wide and
            # short compared to others.
            w = shape.width / 914400.0
            h = shape.height / 914400.0
            if w > 4.0 and h < 2.0:
                # Assign hyperlink
                try:
                    shape.click_action.hyperlink.address = final_url
                except Exception:
                    pass
                break


def _scale_logo_to_fit(image_bytes: bytes, max_width: int, max_height: int) -> bytes:
    """
    Given raw image bytes, return new bytes scaled down to fit within
    ``max_width`` x ``max_height`` EMU.  If the original image is
    smaller than the bounding box it is returned unchanged.  DPI is
    assumed to be 96 for conversion.
    """
    try:
        with Image.open(BytesIO(image_bytes)) as img:
            w, h = img.size
            # Convert max dimensions to pixels
            max_w_px = int(max_width * 96 / 914400)
            max_h_px = int(max_height * 96 / 914400)
            ratio = min(max_w_px / w, max_h_px / h)
            if ratio < 1.0:
                new_size = (int(w * ratio), int(h * ratio))
                img = img.resize(new_size, Image.LANCZOS)
            buf = BytesIO()
            img.save(buf, format='PNG')
            return buf.getvalue()
    except Exception:
        return image_bytes


def create_presentation_from_template(
    topic: str,
    links: List[str],
    final_url: str,
    template_path: str,
    output_pptx: str = 'output.pptx',
    output_pdf: Optional[str] = None,
    app_overrides: Optional[Dict[str, Dict[str, str]]] = None,
    ) -> str:
    """
    Build a PPTX based upon a provided template.  The function
    preserves all original artwork and layout elements while updating
    just the dynamic content.  If more applications are provided
    than there are programme slides in the template, additional
    slides are cloned from the second programme slide to
    accommodate them.

    Parameters
    ----------
    topic: str
        Industry or category displayed on the cover and closing slide.
    links: List[str]
        AppExchange listing URLs; each will occupy one programme slide.
    final_url: str
        URL linked from the closing slide button.
    template_path: str
        Path to the PPTX template file.
    output_pptx: str, optional
        Destination path for the generated PPTX file.
    output_pdf: str, optional
        If provided and LibreOffice is installed the PPTX will be
        converted to PDF at this location.
    app_overrides: dict, optional
        Mapping of URL to a dictionary with keys "name", "developer"
        and optionally ``"logo_path"`` to override extracted data.

    Returns
    -------
    str
        Path to the written PPTX file.
    """
    # Prepare app metadata list
    apps: List[AppMetadata] = []
    overrides = app_overrides or {}
    for link in links:
        link = link.strip()
        meta = None
        if link in overrides:
            ovr = overrides[link]

            # Read logo bytes if provided; if not present we skip
            logo_bytes = None
            logo_mime = 'image/png'
            
            if 'logo_bytes' in ovr and ovr['logo_bytes']:
                logo_bytes = ovr['logo_bytes']
                logo_mime = ovr.get('logo_mime', 'image/png')
            elif 'logo_path' in ovr and ovr['logo_path']:
                try:
                    with open(ovr['logo_path'], 'rb') as f:
                        logo_bytes = f.read()
                        logo_mime = ovr.get('logo_mime', 'image/png')
                except Exception as e:
                    logo_bytes = None
            else:
                print(f"   ⚠️ no logo in overrides")
                
            meta = AppMetadata(
                url=link,
                name=ovr.get('name', ''),
                developer=ovr.get('developer', ''),
                logo_bytes=logo_bytes if logo_bytes else b'',
                logo_mime=logo_mime,
            )
            print(f"   📊 Created AppMetadata: logo_bytes={len(meta.logo_bytes)} bytes")
        else:
            fetched = fetch_app_metadata(link)
            if fetched:
                meta = fetched
        if meta is None:
            # Fallback placeholder metadata if fetching failed
            meta = AppMetadata(
                url=link,
                name='Unknown App',
                developer='Unknown',
                logo_bytes=b'',
                logo_mime='image/png'
            )
        apps.append(meta)
    # Open template
    prs = Presentation(template_path)
    # Determine how many programme slides exist in template; in the
    # supplied file there are 10 (slides 2-11).  We'll treat any
    # slides between the cover (index 0) and closing slide (last) as
    # programme slides.
    total_slides = len(prs.slides)
    closing_index = total_slides - 1
    programme_start = 1
    programme_count = closing_index - programme_start
    needed = len(apps)
    # Clone the second programme slide if more slides are required
    if needed > programme_count:
        extra = needed - programme_count
        for _ in range(extra):
            _clone_slide(prs, programme_start)  # duplicate second programme slide
        # After cloning, programme_count grows accordingly and the
        # closing slide index shifts; update variables
        closing_index = len(prs.slides) - 1
    # Remove unused programme slides if fewer apps are supplied
    if needed < programme_count:
        # remove from the end of programme region until count matches
        remove_count = programme_count - needed
        # Remove slides starting just before closing_index, preserve closing
        for i in range(remove_count):
            _remove_slide(prs, closing_index - 1 - i)
        closing_index = len(prs.slides) - 1
    # Update cover slide
    _update_cover_slide(prs.slides[0], topic)
    # Update programme slides
    for i, app in enumerate(apps):
        slide_index = programme_start + i
        slide = prs.slides[slide_index]
        _update_slide_fields(slide, app, i + 1)
    # Update closing slide
    closing_slide = prs.slides[closing_index]
    _update_closing_slide(closing_slide, topic, final_url)
    
    # Remove comments from specified slides (1, 2, and last slide)
    try:
        last_slide_index = len(prs.slides) - 1  # Dynamically determine last slide
        _remove_comments_from_slides(prs, [0, 1, last_slide_index])  # Slides 1, 2, and last
        print(f"🔍 Removing comments from slides: 1, 2, {last_slide_index + 1} (last)")
    except Exception as e:
        print(f"⚠️ Error removing comments: {e}")

    # Save PPTX
    prs.save(output_pptx)
    # Optionally convert to PDF using LibreOffice
    if output_pdf:
        try:
            subprocess.run(
                [
                    'libreoffice',
                    '--headless',
                    '--convert-to',
                    'pdf',
                    '--outdir',
                    os.path.dirname(os.path.abspath(output_pdf)),
                    os.path.abspath(output_pptx),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
        except Exception:
            # If conversion fails silently ignore
            pass
    return output_pptx


if __name__ == '__main__':  # pragma: no cover
    import argparse
    parser = argparse.ArgumentParser(description='Build a PPTX from a template.')
    parser.add_argument('--topic', required=True, help='Topic or industry for the presentation')
    parser.add_argument('--links', required=True, help='Comma-separated list of AppExchange listing URLs')
    parser.add_argument('--final-url', required=True, help='URL to link from the closing slide')
    parser.add_argument('--template', default='template.pptx', help='Path to the PPTX template')
    parser.add_argument('--output', default='generated.pptx', help='Output PPTX file')
    parser.add_argument('--pdf', default=None, help='Optional output PDF file')
    args = parser.parse_args()
    links = [s.strip() for s in args.links.split(',') if s.strip()]
    create_presentation_from_template(
        topic=args.topic,
        links=links,
        final_url=args.final_url,
        template_path=args.template,
        output_pptx=args.output,
        output_pdf=args.pdf,
    )